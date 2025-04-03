import streamlit as st
import fitz  # PyMuPDF
import pandas as pd
import pytesseract
import textract
from pathlib import Path
import cv2
import numpy as np
from concurrent.futures import ThreadPoolExecutor
from llama_index.core import VectorStoreIndex, Document, Settings
from llama_index.core.node_parser import SentenceSplitter
from llama_index.embeddings.huggingface import HuggingFaceEmbedding
from llama_index.llms.groq import Groq
import tempfile
import shutil
import warnings
import hashlib
from openpyxl import load_workbook
from pptx import Presentation
from docx import Document as DocxDocument
from PIL import Image
from streamlit.components.v1 import html
import time

# Suppress warnings
warnings.filterwarnings("ignore")

# Configuration
max_threads = 10
max_history = 5
supported_exts = [
    ".pdf", ".doc", ".docx", ".xls", ".xlsx", ".ppt", ".pptx",
    ".odt", ".rtf", ".txt", ".csv", ".json", ".html", ".htm",
    ".png", ".jpg", ".jpeg", ".tif", ".tiff", ".bmp"
]

# System setup
@st.cache_resource
def configure_system():
    embed_model = HuggingFaceEmbedding(model_name="baai/bge-small-en")
    groq_llm = Groq(model="llama-3.3-70b-specdec", api_key=st.secrets["k"]["api_key"])

    tesseract_path = shutil.which('tesseract')
    if not tesseract_path:
        raise EnvironmentError('Tesseract not found in path')
    pytesseract.pytesseract.tesseract_cmd = tesseract_path

    return embed_model, groq_llm

embed_model, groq_llm = configure_system()
Settings.embed_model = embed_model
Settings.text_splitter = SentenceSplitter(chunk_size=512, chunk_overlap=50)

# Processing functions
def preprocess_image(image_data):
    img = cv2.imdecode(np.frombuffer(image_data, np.uint8), cv2.IMREAD_COLOR)
    img = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
    _, img = cv2.threshold(img, 0, 255, cv2.THRESH_BINARY | cv2.THRESH_OTSU)
    return cv2.medianBlur(img, 3)

@st.cache_data(max_entries=10, show_spinner=False)
def extract_text_from_image(image_data):
    try:
        processed_img = preprocess_image(image_data)
        return pytesseract.image_to_string(processed_img).strip()
    except Exception as e:
        st.error(f"OCR error: {str(e)}")
        return ""

def process_pdf(file_path):
    text = ""
    with fitz.open(file_path) as doc:
        for page in doc:
            text += page.get_text()
            for img in page.get_images(full=True):
                base_image = doc.extract_image(img[0])
                text += "\n[image]: " + extract_text_from_image(base_image["image"])
    return text

def process_word(file_path):
    text = ""
    try:
        doc = DocxDocument(file_path)
        text = "\n".join([para.text for para in doc.paragraphs])
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                img_data = rel.target_part.blob
                text += "\n[image]: " + extract_text_from_image(img_data)
    except Exception as e:
        st.error(f"Word error: {str(e)}")
    return text

def process_excel(file_path):
    text = ""
    try:
        sheets = pd.ExcelFile(file_path).sheet_names
        for sheet_name in sheets:
            chunks = pd.read_excel(file_path, sheet_name=sheet_name, engine='openpyxl', chunksize=1000)
            sheet_text = f"\nsheet: {sheet_name}\n"
            for chunk in chunks:
                sheet_text += chunk.to_string(index=False) + "\n"
            text += sheet_text

        # Read images if .xlsx
        if file_path.suffix.lower() == ".xlsx":
            wb = load_workbook(file_path, read_only=True)
            for sheet in wb.worksheets:
                for image in sheet._images:
                    text += "\n[image]: " + extract_text_from_image(image._data())
    except Exception as e:
        st.error(f"Excel error: {str(e)}")
    return text

def process_powerpoint(file_path):
    text = ""
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text += shape.text_frame.text + "\n"
                elif shape.shape_type == 13:  # Picture type
                    text += "\n[image]: " + extract_text_from_image(shape.image.blob)
    except Exception as e:
        st.error(f"PPT error: {str(e)}")
    return text

@st.cache_data(max_entries=5, ttl=3600, show_spinner=False)
def process_file(uploaded_file, temp_dir):
    file_path = Path(temp_dir) / uploaded_file.name
    file_path.write_bytes(uploaded_file.getbuffer())
    ext = file_path.suffix.lower()

    try:
        if ext == ".pdf":
            return process_pdf(file_path)
        elif ext in (".docx", ".doc"):
            return process_word(file_path)
        elif ext in (".xlsx", ".xls"):
            return process_excel(file_path)
        elif ext in (".pptx", ".ppt"):
            return process_powerpoint(file_path)
        elif ext in (".png", ".jpg", ".jpeg", ".tif", ".tiff", ".bmp"):
            return extract_text_from_image(file_path.read_bytes())
        else:
            return textract.process(str(file_path)).decode("utf-8")
    except Exception as e:
        st.error(f"Processing error: {str(e)}")
        return ""

# Chat management
def initialize_session():
    if "messages" not in st.session_state:
        st.session_state.messages = []
    if "file_hash" not in st.session_state:
        st.session_state.file_hash = None
    if "query_engine" not in st.session_state:
        st.session_state.query_engine = None

def reset_chat():
    st.session_state.messages = []

def auto_scroll():
    scroll_js = """
    <script>
    setTimeout(() => {
        window.parent.document.querySelectorAll(
            '[data-testid="stVerticalBlock"]'
        ).forEach(el => {
            el.scrollTop = el.scrollHeight;
        });
    }, 100);
    </script>
    """
    html(scroll_js, height=0)

# Main app
def main():
    st.title("📄 Document Analyzer")
    initialize_session()

    # File upload
    uploaded_files = st.file_uploader(
        "Upload documents",
        type=supported_exts,
        accept_multiple_files=True,
        on_change=reset_chat
    )

    # Process files
    if uploaded_files:
        current_hash = hashlib.md5(b''.join(f.getbuffer() for f in uploaded_files)).hexdigest()

        if current_hash != st.session_state.file_hash:  # Check if new files uploaded
            with st.status("Processing files...", expanded=True) as status:
                try:
                    with tempfile.TemporaryDirectory() as temp_dir:
                        with ThreadPoolExecutor(max_threads) as executor:
                            processed = list(executor.map(
                                lambda f: process_file(f, temp_dir), uploaded_files
                            ))

                        processed = "\n\n".join(processed)
                        text= processed.split("\n\n")

                        documents = [Document(text=text) for text in text]
                        st.write(f"Number of documents: {len(documents)} and this is the content: {documents}")
                        Settings.text_splitter = SentenceSplitter(chunk_size=512, chunk_overlap=50)
                        nodes = SentenceSplitter()
                        index = VectorStoreIndex.from_documents(
                            nodes,
                            transformations=[SentenceSplitter(chunk_size=512, chunk_overlap=50)],show_progress=True
                        )
                        return index.as_query_engine(llm=groq_llm)
                     

                        st.session_state.query_engine = VectorStoreIndex.from_documents(
                           nodes , embed_model=embed_model,show_progress=True
                        ).as_query_engine(llm=groq_llm)

                        st.session_state.file_hash = current_hash
                        status.update(label="Processing complete!", state="complete")
                        reset_chat()

                except Exception as e:
                    st.error(f"Processing failed: {str(e)}")
                    reset_chat()

    # Chat interface
    if st.session_state.query_engine:
        st.header("Chat with Documents")

        # Chat container
        chat_container = st.container()
        with chat_container:
            for msg in st.session_state.messages[-max_history:]:
                with st.chat_message(msg["role"]):
                    st.markdown(msg["content"])
            auto_scroll()

        # Input handling
        if prompt := st.chat_input("Ask about your documents"):
            st.session_state.messages.append({"role": "user", "content": prompt})

            with st.spinner("Analyzing..."):
                try:
                    response = st.session_state.query_engine.query(prompt)
                    st.session_state.messages.append({
                        "role": "assistant", 
                        "content": response.response
                    })
                except Exception as e:
                    st.error(f"Query error: {str(e)}")

            st.rerun()

if __name__ == "__main__":
    main()